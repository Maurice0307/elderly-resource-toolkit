from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def add_title_slide(prs, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left, top, width, height = Inches(1), Inches(1.4), Inches(8), Inches(1.2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(17, 62, 120)
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.4), width, Inches(0.9))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(24)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(80, 92, 131)
    footer_box = slide.shapes.add_textbox(left, Inches(6.8), width, Inches(0.5))
    footer_tf = footer_box.text_frame
    footer_tf.text = footer
    footer_tf.paragraphs[0].font.size = Pt(14)
    footer_tf.paragraphs[0].font.color.rgb = RGBColor(140, 140, 140)
    return slide


def add_content_slide(prs, title, bullets, notes=None, emphasize=None, left_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 250, 250)
    if left_color:
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.3), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*left_color)
        shape.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(17, 62, 120)
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(5.2))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body_tf.paragraphs[0]
            p.text = bullet
        else:
            p = body_tf.add_paragraph()
            p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(45, 62, 88)
        p.space_after = Pt(5)
        p.font.bold = i == 0 and emphasize
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes
    return slide


def add_split_slide(prs, title, left_text, right_text, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 128, 0)
    shape.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(17, 62, 120)
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4.2), Inches(4.8))
    left_tf = left_box.text_frame
    left_tf.text = left_text
    left_tf.paragraphs[0].font.size = Pt(20)
    left_tf.paragraphs[0].font.color.rgb = RGBColor(51, 68, 102)
    right_box = slide.shapes.add_textbox(Inches(5.1), Inches(1.6), Inches(4.2), Inches(4.8))
    right_tf = right_box.text_frame
    right_tf.text = right_text
    right_tf.paragraphs[0].font.size = Pt(20)
    right_tf.paragraphs[0].font.color.rgb = RGBColor(51, 68, 102)
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes
    return slide


def main():
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.33)
    add_title_slide(
        prs,
        title='長者資源工具包 ELDERLY TOOLKIT',
        subtitle='以社群賦能，打造高齡照護資源搜尋與政府協力平台',
        footer='2026 / 公私協力提案'
    )
    add_content_slide(
        prs,
        title='問題陳述：高齡化 x 獨居時代的資源斷鏈',
        bullets=[
            '高齡社會加速，獨居長者比例上升，照顧者與志工經常不知道附近可用服務。',
            '現有資源資訊分散在公部門、社團、社區群組與紙本手冊，缺乏統一盤點。',
            '當面臨送醫陪同、交通接送、失智陪伴時，家屬與志工無從快速查詢。',
            '缺乏一個「即時可搜尋、可回報、可同步」的長照資源資料平台。'
        ],
        notes='描述高齡化、獨居和照顧者痛點。強調現況不是只是資料不足，而是資訊分散、無法快速連結可靠服務。',
        left_color=(255, 128, 0)
    )
    add_content_slide(
        prs,
        title='市場機會與社會價值',
        bullets=[
            '台灣高齡化持續攀升，長照服務需求持續擴大。',
            '社區志工與照顧者是最接近用戶的資源發現者。',
            '政府期待「民間協力 + 官方背書」的公共服務補強方案。',
            '開源平台可將「草根資料」轉化為「政府認證的公共資產」。'
        ],
        notes='說明我們如何將市場痛點轉成社會價值，並為政府呈現策略空間。',
        left_color=(0, 102, 204)
    )
    add_split_slide(
        prs,
        title='解法概覽：ELDERLY TOOLKIT',
        left_text='''• 開源論壇式資源庫
• 社群投稿 + 志工回報
• 自動化與人工審查並存
• API 與政府資料同步''',
        right_text='''• 長者友善搜尋 UI
• 一鍵撥號 / 導航 / 分享
• 逐步投稿表單
• 官方認證藍勾機制''',
        notes='展示產品核心價值與使用者兩大面向：社群貢獻與長者可及性。'
    )
    add_content_slide(
        prs,
        title='使用者旅程：志工與家屬如何使用',
        bullets=[
            '1. 打開網站或 LINE LIFF 進入，系統自動定位或輸入地址。',
            '2. 用關鍵字或分類搜尋所需服務，立即看到距離、電話、營業時間。',
            '3. 檢視資源詳頁，確認是否有官方認證、評分與最新回報。',
            '4. 若資訊錯誤可一鍵舉報 / 提出修正建議，平台進行後續審核。'
        ],
        notes='用場景敘述讓政府聽懂「志工/家屬實際怎麼用這個平台」。',
        left_color=(36, 82, 156)
    )
    add_content_slide(
        prs,
        title='目前成果與產品現況',
        bullets=[
            '已完成核心功能：資源 CRUD、區域 moderator、import/export API。',
            'PostgreSQL + Supabase RLS 已實作，具備安全性與角色權限。',
            '支援社群投稿、資源回報、地區審核與批次匯出。',
            '尚待補強：audit logs、輸入驗證、rate limiting、個資合規。'
        ],
        notes='讓政府看到你已經有產品原型，而非空想提案。',
        left_color=(255, 128, 0)
    )
    add_content_slide(
        prs,
        title='平台差異化與競爭優勢',
        bullets=[
            '從「社群資料」到「官方認證」的雙軌審查流程。',
            '針對長者與服務提供者設計的可及性 UI/UX。',
            '政府可直接取用標準 Open Data，降低整合成本。',
            '開源可複製、可延展至其他公共社福領域。'
        ],
        notes='闡述相較於傳統政府單一維護的優勢。',
        left_color=(0, 102, 204)
    )
    add_content_slide(
        prs,
        title='合作模式：路線 B 公私協力',
        bullets=[
            '平台由民間維運，政府提供官方認證與稽核資源。',
            '建立每日/每週資料同步 API 與官方藍勾標示。',
            '試點期間政府派駐 moderator 與平台共享審核策略。',
            '後續可擴展為政府資料捐贈與雙向資料回饋。'
        ],
        notes='清楚描述合作分工，避免政府誤以為要全面接管。',
        left_color=(36, 82, 156)
    )
    add_content_slide(
        prs,
        title='試點計畫與關鍵指標',
        bullets=[
            '試點區域：台北市 / 台中市 / 高雄市（各一行政區）。',
            '試點期程：6 個月。',
            '成功標準：資源更新率 ≥ 80%/月，正確率 ≥ 95%。',
            '社群參與：至少 50 名活躍投稿者與 5 名地區 moderator。'
        ],
        notes='用數字呈現試點可量化成效。',
        left_color=(255, 128, 0)
    )
    add_content_slide(
        prs,
        title='政府需求與下一步',
        bullets=[
            '指派 1 名政府聯絡窗口與 1-2 名 moderator。',
            '授權平台進行 API 測試與資料同步。',
            '簽署 MOU，選定試點區域與驗收指標。',
            '建立每月檢討機制與專案治理。'
        ],
        notes='讓政府知道他們要做什麼，降低提案阻力。',
        left_color=(0, 102, 204)
    )
    add_content_slide(
        prs,
        title='為何現在要做？',
        bullets=[
            '高齡化正進入加速期，資源斷鏈問題更常發生。',
            '政府與社群共同協力，可快速打造全台指南型資源庫。',
            '平台已具備開源架構，未來可複製至身障、社服、兒福等場域。'
        ],
        notes='結尾呼應時機與未來擴展性，強化立即落地感。',
        left_color=(255, 128, 0)
    )
    output = Path('proposal_presentation_pro.pptx')
    prs.save(output)
    print(f'Created {output}')


if __name__ == '__main__':
    main()
