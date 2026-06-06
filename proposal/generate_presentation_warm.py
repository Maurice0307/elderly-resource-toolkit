from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

PRIMARY = RGBColor(214, 85, 0)
SECONDARY = RGBColor(255, 203, 107)
BACKGROUND = RGBColor(255, 245, 230)
TEXT_DARK = RGBColor(74, 43, 11)
TAG_BG = RGBColor(255, 228, 181)
TAG_TEXT = RGBColor(122, 62, 3)


def add_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.3), Inches(2.3))
    header_fill = header.fill
    header_fill.solid()
    header_fill.fore_color.rgb = PRIMARY
    title = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.3), Inches(1.2))
    title_tf = title.text_frame
    title_tf.text = 'ELDERLY TOOLKIT'
    title_tf.paragraphs[0].font.size = Pt(52)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.8))
    sub_tf = subtitle.text_frame
    sub_tf.text = '中高齡資源平台：查詢、學習、分享，陪伴更有溫度'
    sub_tf.paragraphs[0].font.size = Pt(26)
    sub_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    footer = slide.shapes.add_textbox(Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.5))
    footer_tf = footer.text_frame
    footer_tf.text = '2026 / 公私協力提案 • Morris Chiang'
    footer_tf.paragraphs[0].font.size = Pt(14)
    footer_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    return slide


def add_bullet_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    side_fill = side.fill
    side_fill.solid()
    side_fill.fore_color.rgb = PRIMARY
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5))
    body_tf = body.text_frame
    body_tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body_tf.paragraphs[0]
            p.text = bullet
        else:
            p = body_tf.add_paragraph()
            p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)
    if note:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    return slide


def add_tagged_bullet_slide(prs, title, bullets, tags, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.1), Inches(0.1), Inches(13.1), Inches(1.0))
    header_fill = header.fill
    header_fill.solid()
    header_fill.fore_color.rgb = SECONDARY
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    for idx, tag in enumerate(tags):
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 3.2), Inches(1.3), Inches(3.0), Inches(0.45))
        tag_fill = tag_box.fill
        tag_fill.solid()
        tag_fill.fore_color.rgb = TAG_BG
        tag_box.line.color.rgb = SECONDARY
        tag_tf = tag_box.text_frame
        tag_tf.text = tag
        tag_tf.paragraphs[0].font.size = Pt(14)
        tag_tf.paragraphs[0].font.bold = True
        tag_tf.paragraphs[0].font.color.rgb = TAG_TEXT
    body = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.8))
    body_tf = body.text_frame
    body_tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body_tf.paragraphs[0]
            p.text = f'• {bullet}'
        else:
            p = body_tf.add_paragraph()
            p.text = f'• {bullet}'
        p.level = 0
        p.font.size = Pt(21)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(7)
    if note:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    return slide


def add_process_slide(prs, title, steps, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    x = 0.8
    y = 1.4
    width = 2.4
    height = 1.1
    for i, step in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + (width + 0.3) * i), Inches(y), Inches(width), Inches(height))
        box_fill = box.fill
        box_fill.solid()
        box_fill.fore_color.rgb = RGBColor(255, 221, 178)
        box.line.color.rgb = PRIMARY
        tf = box.text_frame
        tf.text = f'{i+1}. {step[0]}'
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        p = tf.add_paragraph()
        p.text = step[1]
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        if i < len(steps) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + (width + 0.3) * i + width), Inches(y + height / 2), Inches(x + (width + 0.3) * i + width + 0.3), Inches(y + height / 2))
            line.line.color.rgb = TEXT_DARK
            line.line.width = Pt(2)
    if note:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    return slide


def add_icon_slide(prs, title, cards, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    x = 0.8
    y = 1.5
    for i, card in enumerate(cards):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + i * 4.1), Inches(y), Inches(1.2), Inches(1.2))
        circle_fill = circle.fill
        circle_fill.solid()
        circle_fill.fore_color.rgb = SECONDARY
        circle.line.color.rgb = PRIMARY
        circle.text_frame.text = card[0]
        circle.text_frame.paragraphs[0].font.size = Pt(24)
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
        text_box = slide.shapes.add_textbox(Inches(x + i * 4.1), Inches(y + 1.4), Inches(3.5), Inches(2.3))
        text_tf = text_box.text_frame
        text_tf.text = card[1]
        text_tf.paragraphs[0].font.size = Pt(18)
        text_tf.paragraphs[0].font.color.rgb = TEXT_DARK
        p = text_tf.add_paragraph()
        p.text = card[2]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(101, 57, 17)
    if note:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    return slide


def add_conclusion_slide(prs, title, columns, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    for i, col in enumerate(columns):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 4.0), Inches(1.4), Inches(3.7), Inches(4.2))
        box_fill = box.fill
        box_fill.solid()
        box_fill.fore_color.rgb = RGBColor(255, 238, 204)
        box.line.color.rgb = PRIMARY
        tf = box.text_frame
        tf.text = col[0]
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        first = True
        for line in col[1:]:
            if first:
                p = tf.add_paragraph()
                first = False
            else:
                p = tf.add_paragraph()
            p.text = f'• {line}'
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_DARK
    if note:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    return slide


def add_two_column_slide(prs, title, left_text, right_text, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.3), Inches(5.8))
    left_tf = left_box.text_frame
    left_tf.text = left_text
    left_tf.paragraphs[0].font.size = Pt(18)
    left_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    right_box = slide.shapes.add_textbox(Inches(6.3), Inches(1.4), Inches(5.3), Inches(5.8))
    right_tf = right_box.text_frame
    right_tf.text = right_text
    right_tf.paragraphs[0].font.size = Pt(18)
    right_tf.paragraphs[0].font.color.rgb = TEXT_DARK
    if note:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    return slide


def main():
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.33)
    add_title_slide(prs)
    add_tagged_bullet_slide(
        prs,
        title='問題陳述：高齡化與資源斷鏈',
        tags=['資源散落', '快速查找難', '長者可及性'],
        bullets=[
            '高齡化與獨居長者增加，照顧者與志工常常不知道身邊已有哪一條可用資源線索。',
            '缺乏系統化資源盤點平台，導致社福人力浪費在「找資源」而非「陪伴長者」。',
            '現有資訊分散在公部門、社團、社區群組與紙本手冊，格式不一且更新不及時。',
            '遇到緊急送醫、交通接送、失智陪伴時，志工和家屬難以快速找到可信賴服務。'
        ],
        note='保留原始痛點描述，讓政府理解高齡服務不是只有長照，而是生活維持與陪伴。'
    )
    add_bullet_slide(
        prs,
        title='我們的解法：中高齡資源平台一頁總覽',
        bullets=[
            '中高齡使用者可查詢附近資源、學習相關知識、分享自己的經驗。',
            '開源資料平台 + 社群貢獻 + 多層審查機制。',
            '可輸出標準化資料（CSV/JSON），支援政府或地方團隊直接匯入。',
            '資源不只是長照，還包含健康、便利、社交、學習與陪伴資源。'
        ],
        note='原有解法內容完整保留，並凸顯平台對政府與社群的價值。'
    )
    add_two_column_slide(
        prs,
        title='使用者流程與 UX 亮點',
        left_text='''主要使用者：長者本人、家屬、志工、社服人員、地方政府職員

設計原則：大字體、高對比、明確 CTA、一鍵通話/導航、語音輔助（TTS）。''',
        right_text='''1. 使用者打開首頁或透過 LINE LIFF 進入，系統自動偵測地點或允許輸入地址。
2. 搜尋關鍵字或分類，結果以卡片式呈現，含距離、電話、營業時間與快速動作。 
3. 點開詳頁查看評價、歷史更新、社群投票與舉報按鈕。 
4. 錯誤資訊可一鍵舉報或提出修正建議，進入社群審核流程。''',
        note='補充 UX 設計與使用流程，讓平台更具人性化。'
    )
    add_bullet_slide(
        prs,
        title='目前成果與短期成效',
        bullets=[
            '已建置核心功能：資源 CRUD、區域 moderator、import/export API。',
            '技術基礎：RLS、角色模型、資料 schema。',
            '現階段支援社群投稿、資源回報、地區審核與批次匯出。',
            '目前短板：缺 audit logs、輸入驗證、rate limiting、個資合規。'
        ],
        note='保留原始成效描述，讓政府明確看到現況與待補強項目。'
    )
    add_bullet_slide(
        prs,
        title='資料品質保證架構',
        bullets=[
            '自動化篩選：格式、重複、惡意連結檢查。',
            '社群投票 + 信任分數機制，先發後審。',
            'Moderator 優先隊列與 SLA 追蹤。',
            '兼顧低門檻參與與高品質輸出。'
        ],
        note='完整呈現原始資料審核與品質機制。'
    )
    add_bullet_slide(
        prs,
        title='政府合作模式：路線 B 建議',
        bullets=[
            '平台由民間維運，政府提供官方認證與稽核人力。',
            '每日/每週資料同步 API + 官方藍勾標示。',
            'MOU → 試點 → 正式採用。',
            '後續可擴展為政府資料捐贈與雙向資料回饋。'
        ],
        note='還原原始合作分工與政府期待。'
    )
    add_bullet_slide(
        prs,
        title='技術需求與合規保證',
        bullets=[
            '必做：Audit Logging、輸入驗證、資料保留政策。',
            '建議：Rate limiting、滲透測試、PIA/隱私影響評估。',
            '輸出格式：OpenAPI + CSV/JSON/XML。',
            '提供 2 週完成 audit logs + dashboard 的時間表。'
        ],
        note='保留技術與合規重點，讓提案更具可執行性。'
    )
    add_bullet_slide(
        prs,
        title='試點計畫與關鍵指標',
        bullets=[
            '推薦試點縣市：台北市、台中市、高雄市（各一行政區）。',
            '期間：6 個月。',
            '成功標準：資源更新率 ≥ 80%／月、資料正確率 ≥ 95%。',
            '社群參與：至少 50 名活躍投稿者與 5 名地區 moderator。'
        ],
        note='完整還原原始試點計畫與驗收指標。'
    )
    add_bullet_slide(
        prs,
        title='成本模型範例',
        bullets=[
            '平台運營（民間）：NTD 100 萬／年。',
            '政府 moderator（5 人）：NTD 300 萬／年。',
            '試點費用（第一年含開發）：NTD 400 萬。',
            '成本遠低於傳統政府單一建置模式。'
        ],
        note='補回成本模型與政府預算參考。'
    )
    add_bullet_slide(
        prs,
        title='風險與緩解策略',
        bullets=[
            '風險：資料品質疑慮、資安責任、預算核定緩慢。',
            '緩解：審計日誌、第三方稽核、分階段試點與 MOU。',
            '強調分階段驗收降低政府風險。'
        ],
        note='這一頁保留原始政府疑慮與對應策略。'
    )
    add_bullet_slide(
        prs,
        title='我們的要求：政府協作項',
        bullets=[
            '指派 1 名政府聯絡窗口與 1-2 位 moderator。',
            '提供試點所需的聯繫管道與行政支援。',
            '批准每日資料同步 API 的存取與測試。',
            '建立每月檢討機制與專案治理。'
        ],
        note='明確列出政府需承擔的事項，降低曖昧。'
    )
    add_bullet_slide(
        prs,
        title='下一步與聯絡方式',
        bullets=[
            '簽署 MOU → 啟動 6 個月試點 → 每月一次檢討會。',
            '快速上線試點，驗證「查詢 + 回報 + 學習」價值。',
            '聯絡人：Morris Chiang（itchiang2025@gmail.com）。'
        ],
        note='呼籲立即啟動與安排首次技術對接。'
    )
    output = Path('proposal_presentation_warm_updated.pptx')
    prs.save(output)
    print(f'Created {output}')


if __name__ == '__main__':
    main()
