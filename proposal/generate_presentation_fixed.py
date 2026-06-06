from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def parse_markdown(md_text: str):
    slides = []
    current = None
    for line in md_text.splitlines():
        stripped = line.strip()
        if stripped.startswith('## '):
            if current:
                slides.append(current)
            title = stripped[3:].strip()
            current = {'title': title, 'bullets': [], 'notes': []}
        elif stripped.startswith('- '):
            if current is not None:
                current['bullets'].append(stripped[2:].strip())
        elif stripped.startswith('**講稿要點：**'):
            note = stripped[len('**講稿要點：**'):].strip()
            if current is not None:
                current['notes'].append(note)
        elif stripped == '---':
            continue
        elif stripped and current is not None:
            if current['notes']:
                current['notes'][-1] += ' ' + stripped
            elif current['bullets']:
                current['bullets'][-1] += ' ' + stripped
    if current:
        slides.append(current)
    return slides


def add_slide(prs, title, bullets, notes):
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(1)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    if bullets:
        body_top = Inches(1.7)
        body_height = Inches(5.5)
        body_shape = slide.shapes.add_textbox(left, body_top, width, body_height)
        body_frame = body_shape.text_frame
        body_frame.word_wrap = True
        body_frame.margin_top = Inches(0)
        body_frame.margin_bottom = Inches(0)
        body_frame.margin_left = Inches(0.1)
        body_frame.margin_right = Inches(0.1)
        first = True
        for bullet in bullets:
            if first:
                p = body_frame.paragraphs[0]
                p.text = bullet
                first = False
            else:
                p = body_frame.add_paragraph()
                p.text = bullet
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
            p.space_after = Pt(6)
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = '\n'.join(notes)
    return slide


def build_pptx(slides, output_path: Path):
    prs = Presentation()
    # remove default slides if any
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
    for slide in slides:
        add_slide(prs, slide['title'], slide['bullets'], slide['notes'])
    prs.save(output_path)


def main():
    root = Path(__file__).parent
    md_path = root / 'proposal_presentation.md'
    output_path = root / 'proposal_presentation_fixed.pptx'
    md_text = md_path.read_text(encoding='utf-8')
    slides = parse_markdown(md_text)
    build_pptx(slides, output_path)
    print(f'Created {output_path}')


if __name__ == '__main__':
    main()
