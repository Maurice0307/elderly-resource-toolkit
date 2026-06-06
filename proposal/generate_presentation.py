from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def parse_markdown(md_text: str):
    slides = []
    current = None
    for line in md_text.splitlines():
        stripped = line.strip()
        if stripped.startswith('## '):
            if current:
                slides.append(current)
            title = stripped[3:].strip()
            current = {'title': title, 'content': [], 'notes': []}
        elif stripped.startswith('- '):
            if current is not None:
                current['content'].append(stripped[2:].strip())
        elif stripped.startswith('**講稿要點：**'):
            note = stripped[len('**講稿要點：**'):].strip()
            if current is not None:
                current['notes'].append(note)
        elif stripped.startswith('**') and stripped.endswith('**') and '講稿要點' in stripped:
            # fallback ignore
            pass
        elif stripped == '---':
            continue
        elif stripped:
            # handle continuation of notes or bullets without dash
            if current is None:
                continue
            if current['notes'] and not stripped.startswith('-'):
                current['notes'][-1] += ' ' + stripped
            elif current['content'] and stripped.startswith('**'):  # maybe note line
                continue
            else:
                pass
    if current:
        slides.append(current)
    return slides


def build_pptx(slides, output_path: Path):
    prs = Presentation()
    # use title only slide layout for first slide
    first = True
    for slide in slides:
        if first:
            slide_layout = prs.slide_layouts[0]
            first = False
        else:
            slide_layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(slide_layout)
        title = s.shapes.title
        title.text = slide['title']
        if slide_layout == prs.slide_layouts[0]:
            if slide['content']:
                subtitle = s.placeholders[1]
                subtitle.text = '\n'.join(slide['content'])
        else:
            body = s.shapes.placeholders[1].text_frame
            body.text = slide['content'][0] if slide['content'] else ''
            for bullet in slide['content'][1:]:
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
        if slide['notes']:
            notes_slide = s.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = '\n'.join(slide['notes'])
    prs.save(output_path)


def main():
    root = Path(__file__).parent
    md_path = root / 'proposal_presentation.md'
    output_path = root / 'proposal_presentation.pptx'
    md_text = md_path.read_text(encoding='utf-8')
    slides = parse_markdown(md_text)
    build_pptx(slides, output_path)
    print(f'Created {output_path}')


if __name__ == '__main__':
    main()
